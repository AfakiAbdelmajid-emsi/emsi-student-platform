from fastapi import Depends, APIRouter, HTTPException
import os
from groq import Groq
from dotenv import load_dotenv
from pydantic import BaseModel
from core.database import supabase
import fitz  # PyMuPDF
import httpx
import re
from io import BytesIO
from pptx import Presentation
from docx import Document
from core.security import get_current_user
from datetime import datetime
from api.AIChat.schemas import ConversationCreate, MessageCreate, ChatMessage
from uuid import uuid4

load_dotenv()

client = Groq(api_key=os.getenv("GROQ_API_KEY"))
router = APIRouter()


@router.post("/start-conversation")
async def start_conversation(payload: dict, user = Depends(get_current_user)):
    try:
        message = payload.get("message", "").strip()
        if not message:
            raise HTTPException(status_code=400, detail="Initial message required to create conversation.")

        title = generate_title_from_message(message)

        new_conversation = {
            "user_id": user["id"],
            "title": title,
            "created_at": datetime.utcnow().isoformat(),
            "updated_at": datetime.utcnow().isoformat()
        }
        result = supabase.table("conversations").insert(new_conversation).execute()
        return {"conversation_id": result.data[0]["id"]}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))



def generate_title_from_message(message: str) -> str:
    message = message.strip()
    return message[:40].capitalize() + ('...' if len(message) > 40 else '')



@router.get("/get-conversations")
async def get_conversations(user = Depends(get_current_user)):
    try:
        result = supabase.table("conversations").select("*").eq("user_id", user["id"]).order("updated_at", desc=True).execute()
        return {"conversations": result.data}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/save-message")
async def save_message(message: MessageCreate):
    try:
        new_message = {
            "id": str(uuid4()),
            "conversation_id": message.conversation_id,
            "role": message.role,
            "content": message.content,
            "created_at": datetime.utcnow().isoformat()
        }

        supabase.table("conversations").update({
            "updated_at": datetime.utcnow().isoformat()
        }).eq("id", message.conversation_id).execute()

        result = supabase.table("messages").insert(new_message).execute()
        return {"success": True}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/get-messages/{conversation_id}")
async def get_messages(conversation_id: str):
    try:
        result = supabase.table("messages").select("*").eq("conversation_id", conversation_id).order("created_at", desc=False).execute()
        return {"messages": result.data}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/ai-chat")
async def ai_chat(messages: list[ChatMessage]):
    try:
        last_message = messages[-1].content.lower()

        try:
            file_name = extract_file_name_from_message(last_message)
            file_result = supabase.table("files").select("file_name").eq("file_name", file_name).execute()
            if file_result.data:
                return await explain_file(file_name)
        except HTTPException as fe:
            if fe.status_code != 400:
                raise fe

        chat_completion = client.chat.completions.create(
            messages=[msg.dict() for msg in messages],
            model="llama-3.3-70b-versatile",
        )
        reply = chat_completion.choices[0].message.content
        return {"reply": reply}
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"AI request failed: {str(e)}")


@router.post("/explain_file")
async def explain_file(file_name: str):
    try:
        file_result = supabase.table("files").select("*").eq("file_name", file_name).execute()
        if not file_result.data:
            raise HTTPException(status_code=404, detail="File not found")

        file_path = file_result.data[0]["file_path"]
        signed_url = supabase.storage.from_("filesb").create_signed_url(file_path, 3600)
        if isinstance(signed_url, dict) and "error" in signed_url:
            raise HTTPException(status_code=400, detail=signed_url["error"])

        file_content = await download_file(signed_url['signedURL'])
        extracted_text = extract_text_by_file_type(file_name, file_content)

        ai_response = await ai_chat([
            ChatMessage(role="user", content=f"Explain this content: {extracted_text}")
        ])
        return {"reply": ai_response['reply']}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to explain the file: {str(e)}")


async def download_file(url: str) -> bytes:
    async with httpx.AsyncClient() as client:
        response = await client.get(url)
        return response.content


def extract_text_by_file_type(file_name: str, file_content: bytes) -> str:
    if file_name.lower().endswith('.pdf'):
        return extract_pdf_text(file_content)
    elif file_name.lower().endswith(('.ppt', '.pptx')):
        return extract_pptx_text(file_content)
    elif file_name.lower().endswith('.docx'):
        return extract_docx_text(file_content)
    elif file_name.lower().endswith('.txt'):
        return file_content.decode('utf-8')
    else:
        raise HTTPException(status_code=400, detail="Unsupported file format")


def extract_pdf_text(pdf_content: bytes) -> str:
    doc = fitz.open(stream=pdf_content, filetype="pdf")
    return "\n".join(page.get_text("text") for page in doc)


def extract_pptx_text(pptx_content: bytes) -> str:
    prs = Presentation(BytesIO(pptx_content))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


def extract_docx_text(docx_content: bytes) -> str:
    doc = Document(BytesIO(docx_content))
    return "\n".join(paragraph.text for paragraph in doc.paragraphs)


def extract_file_name_from_message(message: str) -> str:
    match = re.search(r"([a-zA-Z0-9_\-]+\.(pdf|pptx?|docx|txt))", message)
    if match:
        return match.group(1)
    else:
        raise HTTPException(status_code=400, detail="File name not found in the request")
